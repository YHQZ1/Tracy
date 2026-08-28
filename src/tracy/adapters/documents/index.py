"""Extract, index, and search text from Moodle documents."""

from __future__ import annotations

import hashlib
import json
import re
from collections.abc import Iterable
from dataclasses import asdict
from pathlib import Path

from tracy.domain.entities import Document, DocumentChunk, SyncSnapshot


class DocumentIndexError(RuntimeError):
    """Raised when a document index cannot be loaded or built."""


_STOP_WORDS = {
    "a",
    "about",
    "an",
    "and",
    "are",
    "can",
    "do",
    "for",
    "how",
    "in",
    "is",
    "me",
    "my",
    "of",
    "on",
    "should",
    "the",
    "this",
    "to",
    "what",
    "which",
    "where",
    "with",
    "you",
}
_ROMAN_NUMERALS = {"i": "1", "ii": "2", "iii": "3", "iv": "4", "v": "5"}
_SYLLABUS_MARKERS = {
    "course",
    "evaluation",
    "hours",
    "learning",
    "objectives",
    "outline",
    "pedagogy",
    "prerequisites",
}


def _tokens(value: str) -> list[str]:
    tokens: list[str] = []
    for token in re.findall(r"[a-z0-9]+", value.casefold()):
        token = _ROMAN_NUMERALS.get(token, token)
        if token not in _STOP_WORDS and (len(token) > 1 or token.isdigit()):
            tokens.append(token)
    return tokens


def _chunks(text: str, limit: int = 1400) -> Iterable[str]:
    paragraphs = [" ".join(part.split()) for part in text.splitlines() if part.strip()]
    current = ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 1 <= limit:
            current = f"{current} {paragraph}".strip()
            continue
        if current:
            yield current
        current = paragraph[:limit]
    if current:
        yield current


def _extract_text(path: Path) -> list[tuple[int | None, str]]:
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        try:
            import pymupdf
        except ModuleNotFoundError as error:
            raise DocumentIndexError(
                "PDF support is not installed. Run `uv sync --extra documents`."
            ) from error
        with pymupdf.open(path) as pdf:
            return [(page_number, page.get_text()) for page_number, page in enumerate(pdf, 1)]

    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ModuleNotFoundError as error:
            raise DocumentIndexError(
                "PPTX support is not installed. Run `uv sync --extra documents`."
            ) from error
        presentation = Presentation(path)
        return [
            (
                slide_number,
                "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")),
            )
            for slide_number, slide in enumerate(presentation.slides, 1)
        ]

    if suffix == ".docx":
        try:
            from docx import Document as WordDocument
        except ModuleNotFoundError as error:
            raise DocumentIndexError(
                "DOCX support is not installed. Run `uv sync --extra documents`."
            ) from error
        document = WordDocument(path)
        paragraphs = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            paragraphs.extend(" | ".join(cell.text for cell in row.cells) for row in table.rows)
        return [(None, "\n".join(paragraphs))]

    try:
        return [(None, path.read_text(encoding="utf-8", errors="ignore"))]
    except UnicodeDecodeError:
        return []


def extract_document(document: Document, course_name: str) -> list[DocumentChunk]:
    """Extract searchable chunks from one downloaded document."""

    if not document.local_path or not document.local_path.exists():
        return []
    chunks: list[DocumentChunk] = []
    for page, text in _extract_text(document.local_path):
        for index, chunk_text in enumerate(_chunks(text)):
            digest = hashlib.sha1(
                f"{document.id}:{page}:{index}:{chunk_text}".encode(), usedforsecurity=False
            ).hexdigest()[:16]
            chunks.append(
                DocumentChunk(
                    id=digest,
                    document_id=document.id,
                    course_id=document.course_id,
                    course_name=course_name,
                    document_name=document.name,
                    source_url=document.source_url,
                    text=chunk_text,
                    page=page,
                )
            )
    return chunks


class DocumentIndex:
    """A small deterministic lexical index for the first retrieval slice."""

    def __init__(self, chunks: Iterable[DocumentChunk] = ()) -> None:
        self.chunks = tuple(chunks)

    def search(self, query: str, limit: int = 5) -> list[DocumentChunk]:
        query_tokens = set(_tokens(query))
        if not query_tokens:
            return []

        ranked: list[tuple[int, DocumentChunk]] = []
        for chunk in self.chunks:
            name_tokens = set(_tokens(chunk.document_name))
            course_tokens = set(_tokens(chunk.course_name))
            text_tokens = set(_tokens(chunk.text))
            score = (
                5 * len(query_tokens & name_tokens)
                + 2 * len(query_tokens & course_tokens)
                + 2 * len(query_tokens & text_tokens)
            )
            if "syllabus" in query_tokens:
                score += 3 * len(_SYLLABUS_MARKERS & text_tokens)
            if "lab" in course_tokens and "lab" not in query_tokens:
                score -= 10
            if score:
                ranked.append((score, chunk))
        ranked.sort(key=lambda item: (-item[0], item[1].document_name, item[1].page or 0))
        return [chunk for _, chunk in ranked[:limit]]


class JsonDocumentIndexStore:
    """Persist the lexical index as inspectable JSON."""

    def __init__(self, data_dir: Path) -> None:
        self.path = data_dir / "document-index.json"

    def save(self, index: DocumentIndex) -> None:
        self.path.parent.mkdir(parents=True, exist_ok=True)
        payload = [asdict(chunk) for chunk in index.chunks]
        self.path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    def load(self) -> DocumentIndex:
        if not self.path.exists():
            raise FileNotFoundError(
                f"No document index found at {self.path}. Run `tracy index` first."
            )
        payload = json.loads(self.path.read_text(encoding="utf-8"))
        return DocumentIndex(DocumentChunk(**item) for item in payload)


def build_document_index(snapshot: SyncSnapshot, data_dir: Path) -> DocumentIndex:
    """Extract all downloaded documents and persist their searchable chunks."""

    course_names = {course.id: course.name for course in snapshot.courses}
    chunks = [
        chunk
        for document in snapshot.documents
        for chunk in extract_document(
            document, course_names.get(document.course_id, "Unknown course")
        )
    ]
    index = DocumentIndex(chunks)
    JsonDocumentIndexStore(data_dir).save(index)
    return index
